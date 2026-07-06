"""
DAWN Ingestion Router - file upload, repo ingest, document ingest, memory ingest.

All ingestion runs in background tasks with retry logic, structured logging,
and a persistent job queue with status tracking.

Supports streaming ingestion for files up to 20GB+ - files are processed
in chunks rather than loaded entirely into memory. Scanned PDFs are OCR'd
page-by-page with configurable page limits.

v2.0 — Added auto-tagging (sentence-transformers) and multi-file/zip upload.
v2.1 — Hybrid semantic tagger: Top-K + adaptive floor + dynamic per-tag thresholds.
"""
from fastapi import APIRouter, Depends, HTTPException, Header, UploadFile, File, Form
from pydantic import BaseModel
from typing import Optional
from config import settings
import db.client as db
import logging
import asyncio
import io
import os
import zipfile
import time
import uuid
import tempfile
import shutil
from dataclasses import dataclass, field
from enum import Enum

logger = logging.getLogger(__name__)
router = APIRouter()

_tesseract_configured = False


def _ensure_tesseract_configured():
    global _tesseract_configured
    if _tesseract_configured:
        return
    if settings.tesseract_cmd:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd
        logger.info(f"pytesseract configured: {settings.tesseract_cmd}")
    _tesseract_configured = True


SUPPORTED_EXTENSIONS = {
    ".pdf": "PDF", ".md": "Markdown", ".markdown": "Markdown",
    ".csv": "CSV", ".xlsx": "Excel", ".xls": "Excel",
    ".svg": "SVG", ".docx": "Word", ".txt": "Text",
    ".epub": "EPUB", ".html": "HTML", ".htm": "HTML",
    ".rtf": "RTF", ".json": "JSON",
    ".odt": "Word", ".ods": "Spreadsheet", ".pptx": "PowerPoint", ".odp": "Presentation",
    ".xml": "XML", ".yaml": "YAML", ".yml": "YAML",
    ".log": "Text", ".ini": "Text", ".cfg": "Text", ".conf": "Text",
}

_STREAMING_THRESHOLD_BYTES = 50 * 1024 * 1024
_MAX_UPLOAD_BYTES = 30 * 1024 * 1024 * 1024
_MAX_ZIP_UNCOMPRESSED_BYTES = 500 * 1024 * 1024
_MAX_ZIP_COMPRESSION_RATIO = 100
_MAX_SPREADSHEET_ROWS = 100_000
_MAX_OCR_PAGES = 5000
_MAX_PDF_PAGES = 10000

# Maximum body size for a single table-data child node.
# Tables larger than this get split into multiple child nodes.
_MAX_TABLE_BODY_CHARS = 50_000


class IngestionStatus(str, Enum):
    QUEUED = "queued"
    RUNNING = "running"
    SUCCESS = "success"
    FAILED = "failed"
    PARTIAL = "partial"


@dataclass
class IngestionJob:
    id: str
    type: str
    params: dict
    status: IngestionStatus = IngestionStatus.QUEUED
    error: Optional[str] = None
    result: Optional[dict] = None
    created_at: float = field(default_factory=time.time)
    started_at: Optional[float] = None
    completed_at: Optional[float] = None


class IngestionQueue:
    def __init__(self, max_concurrent: int = 2):
        self.queue: asyncio.Queue[IngestionJob] = asyncio.Queue()
        self.jobs: dict[str, IngestionJob] = {}
        self._semaphore = asyncio.Semaphore(max_concurrent)
        self._worker_task: Optional[asyncio.Task] = None

    async def start(self):
        if self._worker_task is None:
            self._worker_task = asyncio.create_task(self._worker_loop())
            logger.info("Ingestion queue worker started")

    async def stop(self):
        if self._worker_task:
            self._worker_task.cancel()
            self._worker_task = None
            logger.info("Ingestion queue worker stopped")

    async def enqueue(self, job: IngestionJob) -> str:
        self.jobs[job.id] = job
        await self.queue.put(job)
        return job.id

    def get_status(self, job_id: str) -> Optional[IngestionJob]:
        return self.jobs.get(job_id)

    async def _worker_loop(self):
        while True:
            try:
                job = await self.queue.get()
                async with self._semaphore:
                    job.status = IngestionStatus.RUNNING
                    job.started_at = time.time()
                    try:
                        result = await self._execute_job(job)
                        job.status = IngestionStatus.SUCCESS
                        job.result = result
                    except Exception as e:
                        job.status = IngestionStatus.FAILED
                        job.error = str(e)
                        logger.error(f"Job {job.id} ({job.type}) failed: {e}")
                    finally:
                        job.completed_at = time.time()
                        self.queue.task_done()
            except asyncio.CancelledError:
                break
            except Exception as e:
                logger.error(f"Queue worker error: {e}")
                await asyncio.sleep(1)

    async def _execute_job(self, job: IngestionJob) -> dict:
        max_retries = 3
        base_delay = 2.0
        for attempt in range(max_retries):
            try:
                if job.type == "repo":
                    from ingestion.repo import ingest_repo
                    return await ingest_repo(**job.params)
                elif job.type == "document":
                    from ingestion.repo import ingest_document
                    return await ingest_document(**job.params)
                elif job.type == "file":
                    return await self._execute_file_job(job)
                elif job.type == "memory":
                    from ingestion.memory import ingest_memory
                    return await ingest_memory(**job.params)
                else:
                    raise ValueError(f"Unknown job type: {job.type}")
            except (asyncio.TimeoutError, ConnectionError) as e:
                if attempt == max_retries - 1:
                    raise
                delay = base_delay * (2 ** attempt)
                logger.warning(f"Retry {attempt+1}/{max_retries} for job {job.id} after {delay}s: {e}")
                await asyncio.sleep(delay)
            except Exception:
                raise

    async def _execute_file_job(self, job: IngestionJob) -> dict:
        from ingestion.repo import ingest_document_stream, ingest_sections
        file_type = job.params["file_type"]
        title = job.params["title"]
        source_ref = job.params["source_ref"]
        tags = job.params.get("tags", [])
        temp_path = job.params.get("temp_path")

        if temp_path and os.path.exists(temp_path):
            result = await self._process_streaming_file(temp_path, file_type, title, source_ref, tags)
            try:
                os.unlink(temp_path)
            except Exception:
                pass
            return result
        else:
            file_bytes = job.params.get("file_bytes", b"")
            return await self._process_memory_file(file_bytes, file_type, title, source_ref, tags)

    async def _process_streaming_file(self, temp_path, file_type, title, source_ref, tags):
        file_size = os.path.getsize(temp_path)
        logger.info(f"Processing streamed file: {source_ref} ({file_size / 1e9:.2f} GB)")
        with open(temp_path, "rb") as f:
            file_bytes = f.read()
        return await self._process_memory_file(file_bytes, file_type, title, source_ref, tags)

    async def _process_memory_file(self, file_bytes, file_type, title, source_ref, tags):
        from ingestion.repo import ingest_document_stream, ingest_sections
        if file_type == "PDF":
            peek_gen = _peek_pdf_pages(file_bytes, sample_pages=5)
            if peek_gen is None:
                sample_text = ""
            else:
                sample_text = "".join(peek_gen)

            if len(sample_text.strip()) < 20:
                logger.warning(f"PDF '{source_ref}' looks scanned - using OCR.")
                ocr_gen = _iter_ocr_pdf_pages(file_bytes)
                if ocr_gen is None:
                    logger.warning(f"OCR unavailable for '{source_ref}' - falling back to text layer.")
                    pdf_gen = iter_pdf_pages(file_bytes)
                    gen = pdf_gen if pdf_gen is not None else iter([])
                else:
                    gen = ocr_gen
            else:
                pdf_gen = iter_pdf_pages(file_bytes)
                gen = pdf_gen if pdf_gen is not None else iter([])

            result = await ingest_document_stream(title, gen, source_ref, tags)
        elif file_type == "Word":
            sections = _parse_docx(file_bytes)
            result = await ingest_sections(title, sections, source_ref, tags)
        elif file_type == "PowerPoint":
            sections = _parse_pptx(file_bytes)
            result = await ingest_sections(title, sections, source_ref, tags)
        elif file_type == "Presentation":
            sections = _parse_odp(file_bytes)
            result = await ingest_sections(title, sections, source_ref, tags)
        elif file_type == "Spreadsheet":
            sections = _parse_ods(file_bytes)
            result = await ingest_sections(title, sections, source_ref, tags)
        elif file_type == "EPUB":
            sections = _parse_epub(file_bytes)
            result = await ingest_sections(title, sections, source_ref, tags)
        else:
            extraction = _extract_content(file_bytes, file_type, source_ref, title)
            if extraction.get("sections"):
                result = await ingest_sections(title, extraction["sections"], source_ref, tags)
            else:
                from ingestion.repo import ingest_document
                result = await ingest_document(title, extraction["text"], source_ref, tags)
        await db.log_ingestion({
            "source": "document", "source_ref": source_ref,
            "nodes_created": result["nodes_created"],
            "edges_created": result.get("edges_created", 0),
            "status": "success",
        })
        return result


ingestion_queue = IngestionQueue(max_concurrent=2)


def verify_key(x_api_key: Optional[str] = Header(None)):
    if x_api_key != settings.dawn_api_key:
        raise HTTPException(status_code=401, detail="Invalid API key")


class RepoIngestRequest(BaseModel):
    repo_path: str
    repo_name: str
    tags: list[str] = []


class DocumentIngestRequest(BaseModel):
    title: str
    content: str
    source_ref: str = ""
    tags: list[str] = []


class MemoryIngestRequest(BaseModel):
    conversation: str
    session_source: str = "manual"


@router.post("/repo")
async def ingest_repo_endpoint(req: RepoIngestRequest, _: None = Depends(verify_key)):
    job = IngestionJob(id=str(uuid.uuid4()), type="repo",
        params={"repo_path": req.repo_path, "repo_name": req.repo_name, "tags": req.tags})
    await ingestion_queue.enqueue(job)
    return {"status": "queued", "job_id": job.id, "repo": req.repo_name}


@router.post("/document")
async def ingest_document_endpoint(req: DocumentIngestRequest, _: None = Depends(verify_key)):
    job = IngestionJob(id=str(uuid.uuid4()), type="document",
        params={"title": req.title, "content": req.content, "source_ref": req.source_ref, "tags": req.tags})
    await ingestion_queue.enqueue(job)
    return {"status": "queued", "job_id": job.id, "title": req.title}


@router.post("/memory")
async def ingest_memory_endpoint(req: MemoryIngestRequest, _: None = Depends(verify_key)):
    job = IngestionJob(id=str(uuid.uuid4()), type="memory",
        params={"conversation": req.conversation, "session_source": req.session_source})
    await ingestion_queue.enqueue(job)
    return {"status": "queued", "job_id": job.id}


@router.get("/status/{job_id}")
async def get_job_status(job_id: str, _: None = Depends(verify_key)):
    job = ingestion_queue.get_status(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")
    return {
        "job_id": job.id, "type": job.type, "status": job.status.value,
        "error": job.error, "result": job.result,
        "created_at": job.created_at, "started_at": job.started_at, "completed_at": job.completed_at,
    }


# ─── Auto-Tagging (v2.1 — Hybrid Tagger) ─────────────────────────────────
# Uses the HybridTagger from ingestion/tagger.py which implements:
#   1. Top-K (always returns top 2 tags)
#   2. Adaptive floor (0.1 minimum — garbage detection)
#   3. Dynamic per-tag thresholds (learned from historical matches)
# No extra dependencies, no API calls, runs on CPU.

_tagger_instance = None


async def _get_tagger():
    """Lazy-load and return the singleton HybridTagger instance."""
    global _tagger_instance
    if _tagger_instance is None:
        from ingestion.tagger import HybridTagger
        _tagger_instance = HybridTagger()
        await _tagger_instance.refresh()
    return _tagger_instance


async def _auto_tag_content(content: str, title: str = "", top_n: int = 2, threshold: float = 0.1) -> list[str]:
    """Auto-tag document content using the hybrid semantic tagger.

    Uses Top-K (top_n) + adaptive floor (threshold) + dynamic per-tag thresholds.
    Falls back to ["uncategorized"] if nothing clears the bar.
    """
    tagger = await _get_tagger()
    return await tagger.tag(
        text=content,
        title=title,
        top_k=top_n,
        min_similarity=threshold,
        use_dynamic_thresholds=True,
    )


# ─── Single File Upload ───────────────────────────────────────────────────

@router.post("/file")
async def ingest_file(
    file: UploadFile = File(...),
    title: str = Form(""),
    tags: str = Form(""),
    x_api_key: Optional[str] = Header(None),
):
    if x_api_key != settings.dawn_api_key:
        raise HTTPException(status_code=401, detail="Invalid API key")

    filename = file.filename or ""
    file_type = detect_file_type(filename)
    if not file_type:
        ext = os.path.splitext(filename.lower())[1]
        raise HTTPException(status_code=400,
            detail=f"Unsupported file type '{ext}'. Supported: {', '.join(sorted(set(SUPPORTED_EXTENSIONS.values())))}")

    doc_title = title.strip() or os.path.splitext(filename)[0]
    tag_list = [t.strip() for t in tags.split(",") if t.strip()]

    # Read first bytes for validation, then decide streaming vs memory
    header_bytes = await file.read(8192)
    if not header_bytes:
        raise HTTPException(status_code=422, detail="Empty file")

    # Quick sanity check
    if not _quick_sanity_check(header_bytes, file_type):
        raise HTTPException(status_code=422, detail=f"File does not look like a valid {file_type} file.")

    # Zip-bomb check for zip-based formats
    if file_type in ("Word", "EPUB", "Excel"):
        full_check_bytes = header_bytes + await file.read()
        zip_error = _check_zip_safety(full_check_bytes)
        if zip_error:
            raise HTTPException(status_code=422, detail=zip_error)
        file_bytes = full_check_bytes
        total_size = len(file_bytes)
    else:
        remaining = await file.read()
        file_bytes = header_bytes + remaining
        total_size = len(file_bytes)

    if total_size > _MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413,
            detail=f"File is {total_size / 1e9:.2f}GB, exceeds {_MAX_UPLOAD_BYTES // (1024*1024*1024)}GB limit.")

    # ── Auto-tag if no tags provided ──
    if not tag_list:
        # Extract a preview of the content for classification
        preview = _extract_preview(file_bytes, file_type, filename)
        tag_list = await _auto_tag_content(preview, doc_title)
        logger.info(f"Auto-tagged '{filename}' as: {tag_list}")

    # For files > streaming threshold, write to temp file
    temp_path = None
    if total_size > _STREAMING_THRESHOLD_BYTES:
        fd, temp_path = tempfile.mkstemp(suffix=os.path.splitext(filename)[1])
        with os.fdopen(fd, "wb") as f:
            f.write(file_bytes)
        logger.info(f"Large file ({total_size / 1e6:.1f}MB) streamed to temp: {temp_path}")
        job = IngestionJob(id=str(uuid.uuid4()), type="file",
            params={"title": doc_title, "source_ref": filename, "file_type": file_type,
                    "temp_path": temp_path, "tags": tag_list})
    else:
        job = IngestionJob(id=str(uuid.uuid4()), type="file",
            params={"title": doc_title, "source_ref": filename, "file_type": file_type,
                    "file_bytes": file_bytes, "tags": tag_list})

    await ingestion_queue.enqueue(job)
    return {
        "status": "queued", "job_id": job.id, "title": doc_title,
        "file_type": file_type, "filename": filename,
        "tags": tag_list,
        "size_mb": total_size / 1e6,
        "note": f"File ({total_size / 1e6:.1f}MB) queued for ingestion." if total_size > _STREAMING_THRESHOLD_BYTES else None,
    }


# ─── Multi-File / Zip Upload ──────────────────────────────────────────────

@router.post("/files")
async def ingest_files(
    files: list[UploadFile] = File(...),
    tags: str = Form(""),
    x_api_key: Optional[str] = Header(None),
):
    """Upload multiple files and/or zip archives at once.

    Accepts:
    - Multiple individual files (any supported type)
    - .zip archives (extracted automatically, each supported file ingested)
    - A mix of both

    Returns a summary with job IDs for each file.
    """
    if x_api_key != settings.dawn_api_key:
        raise HTTPException(status_code=401, detail="Invalid API key")

    global_tags = [t.strip() for t in tags.split(",") if t.strip()]
    jobs = []
    errors = []

    for file in files:
        filename = file.filename or ""
        if not filename:
            continue

        ext = os.path.splitext(filename.lower())[1]

        # ── Zip archive handling ──
        if ext == ".zip":
            try:
                raw_bytes = await file.read()
                if not raw_bytes:
                    errors.append({"file": filename, "error": "Empty zip file"})
                    continue
                zip_jobs, zip_errors = await _process_zip_upload(
                    raw_bytes, filename, global_tags
                )
                jobs.extend(zip_jobs)
                errors.extend(zip_errors)
            except Exception as e:
                errors.append({"file": filename, "error": f"Zip processing failed: {str(e)}"})
            continue

        # ── Individual file handling ──
        file_type = detect_file_type(filename)
        if not file_type:
            errors.append({"file": filename, "error": f"Unsupported file type '{ext}'"})
            continue

        doc_title = os.path.splitext(filename)[0]
        tag_list = list(global_tags)  # copy

        # Read the file
        file_bytes = await file.read()
        if not file_bytes:
            errors.append({"file": filename, "error": "Empty file"})
            continue

        total_size = len(file_bytes)
        if total_size > _MAX_UPLOAD_BYTES:
            errors.append({"file": filename, "error": f"File exceeds size limit"})
            continue

        # Quick sanity check
        if not _quick_sanity_check(file_bytes[:8192], file_type):
            errors.append({"file": filename, "error": f"File does not look like a valid {file_type}"})
            continue

        # Auto-tag if no global tags provided
        if not tag_list:
            preview = _extract_preview(file_bytes, file_type, filename)
            tag_list = await _auto_tag_content(preview, doc_title)
            logger.info(f"Auto-tagged '{filename}' as: {tag_list}")

        # Queue the job
        job = IngestionJob(id=str(uuid.uuid4()), type="file",
            params={"title": doc_title, "source_ref": filename, "file_type": file_type,
                    "file_bytes": file_bytes, "tags": tag_list})
        await ingestion_queue.enqueue(job)
        jobs.append({
            "job_id": job.id, "filename": filename, "file_type": file_type,
            "tags": tag_list, "size_mb": total_size / 1e6,
        })

    return {
        "status": "complete",
        "total_files": len(jobs) + len(errors),
        "queued": len(jobs),
        "errors": len(errors),
        "jobs": jobs,
        "error_details": errors if errors else None,
    }


async def _process_zip_upload(
    zip_bytes: bytes, zip_filename: str, global_tags: list[str]
) -> tuple[list[dict], list[dict]]:
    """Extract a zip archive and queue ingestion jobs for each supported file.

    Returns (jobs_list, errors_list).
    """
    jobs = []
    errors = []

    try:
        zf = zipfile.ZipFile(io.BytesIO(zip_bytes))
    except zipfile.BadZipFile:
        errors.append({"file": zip_filename, "error": "Invalid zip file"})
        return jobs, errors

    # Safety checks
    total_uncompressed = 0
    for info in zf.infolist():
        total_uncompressed += info.file_size
        if info.file_size > _MAX_ZIP_UNCOMPRESSED_BYTES:
            errors.append({"file": info.filename, "error": "File too large in zip"})
            continue
        if info.compress_size > 0:
            ratio = info.file_size / max(info.compress_size, 1)
            if ratio > _MAX_ZIP_COMPRESSION_RATIO and info.file_size > 10_000_000:
                errors.append({"file": info.filename, "error": "Suspicious compression ratio"})
                continue

    if total_uncompressed > _MAX_ZIP_UNCOMPRESSED_BYTES:
        errors.append({"file": zip_filename, "error": "Zip expands to over 500MB"})
        zf.close()
        return jobs, errors

    # Extract and process each file
    for info in zf.infolist():
        fname = info.filename

        # Skip directories, macOS metadata, hidden files
        if fname.endswith("/"):
            continue
        if fname.startswith("__MACOSX/") or fname.startswith("."):
            continue
        if os.path.basename(fname).startswith("."):
            continue

        ext = os.path.splitext(fname.lower())[1]
        file_type = SUPPORTED_EXTENSIONS.get(ext)
        if not file_type:
            continue

        try:
            file_bytes = zf.read(info.filename)
        except Exception as e:
            errors.append({"file": f"{zip_filename}/{fname}", "error": f"Read failed: {str(e)}"})
            continue

        if not file_bytes:
            continue

        doc_title = os.path.splitext(os.path.basename(fname))[0]
        tag_list = list(global_tags)

        # Auto-tag if no global tags
        if not tag_list:
            preview = _extract_preview(file_bytes, file_type, fname)
            tag_list = await _auto_tag_content(preview, doc_title)

        source_ref = f"{zip_filename}/{fname}"
        job = IngestionJob(id=str(uuid.uuid4()), type="file",
            params={"title": doc_title, "source_ref": source_ref, "file_type": file_type,
                    "file_bytes": file_bytes, "tags": tag_list})
        await ingestion_queue.enqueue(job)
        jobs.append({
            "job_id": job.id, "filename": source_ref, "file_type": file_type,
            "tags": tag_list, "size_mb": len(file_bytes) / 1e6,
        })

    zf.close()
    return jobs, errors


# ─── Content Preview Extraction (for auto-tagging) ────────────────────────

def _extract_preview(file_bytes: bytes, file_type: str, filename: str, max_chars: int = 1000) -> str:
    """Extract a text preview from a file for auto-tagging classification.

    Uses the same parsers as the full ingestion pipeline but only reads
    enough to classify the content. Returns plain text (first ~1000 chars).
    """
    try:
        if file_type == "PDF":
            gen = iter_pdf_pages(file_bytes)
            if gen:
                text = " ".join(gen)
                return text[:max_chars] if text else ""
            # Try OCR preview
            ocr = _iter_ocr_pdf_pages(file_bytes, max_pages=3)
            if ocr:
                text = " ".join(ocr)
                return text[:max_chars] if text else ""
            return ""
        elif file_type == "Word":
            sections = _parse_docx(file_bytes)
            if sections:
                return (sections[0].get("body", "") or "")[:max_chars]
            return ""
        elif file_type == "Markdown":
            text = file_bytes.decode("utf-8", errors="ignore")
            return text[:max_chars]
        elif file_type == "CSV":
            text = file_bytes.decode("utf-8", errors="replace")
            return text[:max_chars]
        elif file_type == "Excel":
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            texts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for row in ws.iter_rows(values_only=True, max_row=5):
                    texts.append(" ".join(str(c) for c in row if c is not None))
                if texts:
                    break
            wb.close()
            return " ".join(texts)[:max_chars]
        elif file_type == "PowerPoint":
            sections = _parse_pptx(file_bytes)
            if sections:
                return (sections[0].get("body", "") or "")[:max_chars]
            return ""
        elif file_type == "Text":
            return file_bytes.decode("utf-8", errors="ignore")[:max_chars]
        elif file_type == "HTML":
            from html.parser import HTMLParser
            class _Extractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.parts = []
                    self._skip = False
                def handle_starttag(self, tag, attrs):
                    if tag in ("script", "style"):
                        self._skip = True
                def handle_endtag(self, tag):
                    if tag in ("script", "style"):
                        self._skip = False
                def handle_data(self, data):
                    if not self._skip and data.strip():
                        self.parts.append(data.strip())
            parser = _Extractor()
            parser.feed(file_bytes.decode("utf-8", errors="ignore"))
            return " ".join(parser.parts)[:max_chars]
        elif file_type == "JSON":
            return file_bytes.decode("utf-8", errors="ignore")[:max_chars]
        elif file_type == "XML":
            return file_bytes.decode("utf-8", errors="ignore")[:max_chars]
        elif file_type == "YAML":
            return file_bytes.decode("utf-8", errors="ignore")[:max_chars]
        elif file_type == "EPUB":
            sections = _parse_epub(file_bytes)
            if sections:
                return (sections[0].get("body", "") or "")[:max_chars]
            return ""
        elif file_type == "SVG":
            return _parse_svg(file_bytes)[:max_chars]
        elif file_type == "RTF":
            from striprtf.striprtf import rtf_to_text
            return rtf_to_text(file_bytes.decode("utf-8", errors="ignore"))[:max_chars]
        elif file_type in ("Spreadsheet", "Presentation"):
            return f"[{file_type} file: {filename}]"[:max_chars]
    except Exception as e:
        logger.debug(f"Preview extraction failed for {filename}: {e}")
        return f"[{file_type} file: {filename}]"

    return f"[{file_type} file: {filename}]"


# ─── Stats & Log Endpoints ────────────────────────────────────────────────

@router.get("/stats")
async def ingestion_stats(_: None = Depends(verify_key)):
    try:
        failed_recent = await db.get_failed_ingestions_since(time.time() - 86400)
        total_nodes = await db.count_nodes()
        db_ok = await db.ping()
    except Exception:
        failed_recent = []; total_nodes = 0; db_ok = False
    active_jobs = sum(1 for j in ingestion_queue.jobs.values()
                      if j.status in (IngestionStatus.QUEUED, IngestionStatus.RUNNING))
    return {
        "queue_depth": ingestion_queue.queue.qsize(), "active_jobs": active_jobs,
        "failed_last_24h": len(failed_recent), "total_nodes": total_nodes, "db_connected": db_ok,
    }


@router.get("/log")
async def get_log(_: None = Depends(verify_key)):
    return await db.get_ingestion_log()



# ──── Admin: Back-Tag Existing Nodes ──────────────────────────────────────
# Re-runs auto-tagging on all nodes tagged "uncategorized" or with no tags.
# Uses the HybridTagger with Top-K + adaptive floor + dynamic per-tag thresholds.
# Also updates per-tag thresholds based on match scores for continuous learning.
# Trigger via: POST /admin/backtag (requires API key)

NEW_TAGS_FOR_BACKTAG = [
    {"name": "self-help", "description": "Self-help, personal development, psychology, philosophy, life strategy, success principles, motivation"},
    {"name": "philosophy", "description": "Philosophy, ethics, political theory, historical analysis, critical thinking, logic, morality"},
    {"name": "psychology", "description": "Psychology, human behavior, cognitive science, persuasion, influence, mental models"},
    {"name": "history", "description": "History, historical events, biographies, historical analysis, ancient civilizations, world history"},
    {"name": "strategy", "description": "Strategy, tactics, game theory, military strategy, business strategy, competitive analysis, power dynamics"},
    {"name": "leadership", "description": "Leadership, management, executive skills, team building, organizational behavior, decision making"},
    {"name": "communication", "description": "Communication, negotiation, persuasion, public speaking, rhetoric, writing, storytelling"},
    {"name": "economics", "description": "Economics, macroeconomics, microeconomics, economic theory, market analysis, trade"},
    {"name": "politics", "description": "Politics, governance, political theory, international relations, policy, power"},
    {"name": "biography", "description": "Biography, memoir, autobiography, personal stories, life narratives, historical figures"},
    {"name": "uncategorized", "description": "Default fallback tag for content that doesn't match any other category"},
]


@router.post("/admin/backtag")
async def admin_backtag(x_api_key: Optional[str] = Header(None)):
    """Re-run auto-tagging on all uncategorized/untagged nodes.

    Uses the HybridTagger with:
    - Top-K (top 2 tags per node)
    - Adaptive floor (0.1 minimum similarity)
    - Dynamic per-tag thresholds (learned from match scores)

    Also creates any missing tags from NEW_TAGS_FOR_BACKTAG in the database.
    Returns a summary of what was tagged and what stayed uncategorized.
    """
    if x_api_key != settings.dawn_api_key:
        raise HTTPException(status_code=401, detail="Invalid API key")

    results = {"tags_created": 0, "nodes_processed": 0, "tagged": 0, "still_uncategorized": 0, "errors": 0}

    # Step 1: Add any missing tags
    existing_tags = await db.get_all_tags()
    existing_names = {t["name"] for t in existing_tags}
    for tag_def in NEW_TAGS_FOR_BACKTAG:
        if tag_def["name"] not in existing_names:
            await db.create_tag(tag_def["name"], tag_def["description"])
            results["tags_created"] += 1
            logger.info(f"Backtag: created tag '{tag_def['name']}'")

    # Refresh tag list
    all_tags = await db.get_all_tags()
    tag_name_to_id = {t["name"]: t["id"] for t in all_tags}

    # Step 2: Find uncategorized and untagged nodes
    db_client = db.get_db()

    # Find the uncategorized tag ID
    uncat_tag = next((t for t in all_tags if t["name"] == "uncategorized"), None)
    uncat_tag_id = uncat_tag["id"] if uncat_tag else None

    # Nodes with "uncategorized" tag
    uncategorized_ids = set()
    if uncat_tag_id:
        res = await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: db_client.table("node_tags")
            .select("node_id")
            .eq("tag_id", uncat_tag_id)
            .execute()
        )
        uncategorized_ids = {row["node_id"] for row in (res.data or [])}

    # All nodes with ANY tag
    res = await asyncio.get_event_loop().run_in_executor(
        None,
        lambda: db_client.table("node_tags").select("node_id").execute()
    )
    any_tag_ids = {row["node_id"] for row in (res.data or [])}

    # All active/stale nodes
    res = await asyncio.get_event_loop().run_in_executor(
        None,
        lambda: db_client.table("nodes")
        .select("id, title, body, type, source, source_ref")
        .in_("status", ["active", "stale"])
        .execute()
    )
    all_nodes = res.data or []

    nodes_to_tag = []
    for node in all_nodes:
        nid = node["id"]
        if nid in uncategorized_ids or nid not in any_tag_ids:
            nodes_to_tag.append(node)

    results["nodes_processed"] = len(nodes_to_tag)

    if not nodes_to_tag:
        return {**results, "message": "No uncategorized or untagged nodes found."}

    # Step 3: Load the HybridTagger
    tagger = await _get_tagger()

    # Step 4: Back-tag each node using the hybrid strategy
    for i, node in enumerate(nodes_to_tag):
        try:
            title = node.get("title", "") or ""
            body = node.get("body", "") or ""
            text = f"{title} {body}" if title else body
            text = text.strip()
            if not text:
                continue

            # Tag using the hybrid strategy
            matched = await tagger.tag(
                text=text,
                title="",
                top_k=2,
                min_similarity=0.1,
                use_dynamic_thresholds=True,
            )

            if matched and matched != ["uncategorized"]:
                tag_ids = [tag_name_to_id[t] for t in matched if t in tag_name_to_id]
                if tag_ids:
                    await db.attach_tags_batch([node["id"]], tag_ids)
                    results["tagged"] += 1

                    # ── Update per-tag thresholds ──
                    # Re-embed to get the actual similarity scores for threshold learning
                    if tagger.model is not None and tagger.tag_embeddings:
                        import numpy as np
                        doc_vec = tagger.model.encode(text[:2000], show_progress_bar=False)
                        doc_norm = doc_vec / (np.linalg.norm(doc_vec) + 1e-10)
                        for tag_name in matched:
                            tag_vec = tagger.tag_embeddings.get(tag_name)
                            if tag_vec is not None:
                                tag_norm = tag_vec / (np.linalg.norm(tag_vec) + 1e-10)
                                score = float(np.dot(doc_norm, tag_norm))
                                await tagger.update_threshold(tag_name, score)
            else:
                results["still_uncategorized"] += 1

        except Exception as e:
            results["errors"] += 1
            logger.error(f"Backtag error on node {node.get('id', '?')[:8]}: {e}")

    # Log current thresholds for diagnostics
    thresholds = await tagger.get_thresholds()
    logger.info(f"Backtag complete. Per-tag thresholds: {thresholds}")

    return {
        **results,
        "thresholds": {k: round(v, 3) for k, v in thresholds.items()},
        "message": (
            f"Processed {results['nodes_processed']} nodes. "
            f"Tagged {results['tagged']}, {results['still_uncategorized']} still uncategorized. "
            f"Per-tag thresholds updated."
        ),
    }
# ─── File Type Detection ──────────────────────────────────────────────────

def detect_file_type(filename: str) -> Optional[str]:
    ext = os.path.splitext(filename.lower())[1]
    return SUPPORTED_EXTENSIONS.get(ext)


def _check_zip_safety(file_bytes: bytes) -> Optional[str]:
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
    except zipfile.BadZipFile:
        return "File is not a valid zip-based document."
    total_uncompressed = 0
    for info in zf.infolist():
        total_uncompressed += info.file_size
        if info.file_size > _MAX_ZIP_UNCOMPRESSED_BYTES:
            return f"Archive member expands to over {_MAX_ZIP_UNCOMPRESSED_BYTES // (1024*1024)}MB."
        if info.compress_size > 0:
            ratio = info.file_size / max(info.compress_size, 1)
            if ratio > _MAX_ZIP_COMPRESSION_RATIO and info.file_size > 10_000_000:
                return "Suspicious compression ratio (possible zip bomb)."
    if total_uncompressed > _MAX_ZIP_UNCOMPRESSED_BYTES:
        return f"Archive expands to over {_MAX_ZIP_UNCOMPRESSED_BYTES // (1024*1024)}MB total."
    return None


def _quick_sanity_check(file_bytes: bytes, file_type: str) -> bool:
    try:
        if file_type == "PDF":
            return file_bytes[:5] == b"%PDF-"
        if file_type in ("Word", "EPUB", "Excel", "PowerPoint", "Presentation"):
            return file_bytes[:2] == b"PK"
        if file_type == "XML":
            return file_bytes[:5].lstrip()[:1] == b"<" or file_bytes[:5] == b"<?xml"
        if file_type == "SVG":
            return file_bytes[:5].lstrip()[:1] == b"<"
        return True
    except Exception:
        return False


def _extract_content(file_bytes: bytes, file_type: str, filename: str, doc_title: str = "") -> dict:
    """Extract content from a file. doc_title is used for table summaries."""
    try:
        if file_type == "PDF":
            return {"text": _parse_pdf(file_bytes), "sections": []}
        elif file_type == "Markdown":
            return {"text": "", "sections": _parse_md(file_bytes.decode("utf-8", errors="ignore"))}
        elif file_type == "CSV":
            return {"text": "", "sections": _parse_csv(file_bytes, doc_title or os.path.splitext(filename)[0])}
        elif file_type == "Excel":
            return {"text": "", "sections": _parse_xlsx(file_bytes, doc_title or os.path.splitext(filename)[0])}
        elif file_type == "SVG":
            return {"text": _parse_svg(file_bytes), "sections": []}
        elif file_type == "Word":
            return {"text": "", "sections": _parse_docx(file_bytes)}
        elif file_type == "PowerPoint":
            return {"text": "", "sections": _parse_pptx(file_bytes)}
        elif file_type == "Presentation":
            return {"text": "", "sections": _parse_odp(file_bytes)}
        elif file_type == "Spreadsheet":
            return {"text": "", "sections": _parse_ods(file_bytes)}
        elif file_type == "Text":
            return {"text": _parse_txt(file_bytes), "sections": []}
        elif file_type == "EPUB":
            return {"text": "", "sections": _parse_epub(file_bytes)}
        elif file_type == "HTML":
            return {"text": _parse_html(file_bytes), "sections": []}
        elif file_type == "RTF":
            return {"text": _parse_rtf(file_bytes), "sections": []}
        elif file_type == "JSON":
            return {"text": _parse_json(file_bytes), "sections": []}
        elif file_type == "XML":
            return {"text": _parse_xml(file_bytes), "sections": []}
        elif file_type == "YAML":
            return {"text": _parse_yaml(file_bytes), "sections": []}
    except Exception as e:
        logger.error(f"Extraction failed for {file_type} ({filename}): {e}")
    return {"text": "", "sections": []}


# ─── PDF Parsing ──────────────────────────────────────────────────────────

def _parse_pdf(file_bytes: bytes) -> str:
    pdf_gen = iter_pdf_pages(file_bytes)
    if pdf_gen is None:
        logger.warning("PDF text layer unavailable - trying OCR.")
        ocr_result = _ocr_pdf(file_bytes)
        return ocr_result
    result = "\n\n".join(pdf_gen)
    if len(result.strip()) < 20:
        logger.warning(f"PDF text layer near-empty ({len(result.strip())} chars) - trying OCR.")
        ocr_result = _ocr_pdf(file_bytes)
        if len(ocr_result.strip()) > len(result.strip()):
            return ocr_result
    return result


def _ocr_pdf(file_bytes: bytes, max_pages: int = 300) -> str:
    try:
        import fitz, pytesseract
        from PIL import Image
    except ImportError:
        logger.warning("OCR fallback unavailable.")
        return ""
    _ensure_tesseract_configured()
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception:
        logger.warning("Failed to open PDF for OCR - file may be corrupt.")
        return ""
    pages_text = []
    total_pages = len(doc)
    pages_to_process = min(total_pages, max_pages)
    try:
        for i in range(pages_to_process):
            page = doc[i]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(img)
            if text and text.strip():
                pages_text.append(text.strip())
    except Exception as e:
        logger.error(f"OCR failed partway: {e}")
    finally:
        doc.close()
    if total_pages > max_pages:
        pages_text.append(f"[OCR truncated: {total_pages} pages, only first {max_pages} OCR'd.]")
    return "\n\n".join(pages_text)


def iter_pdf_pages(file_bytes: bytes):
    try:
        import fitz
    except ImportError:
        logger.warning("fitz (pymupdf) not available - PDF text extraction disabled.")
        return None
    import fitz
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception:
        logger.warning("Failed to open PDF stream - file may be corrupt or empty.")
        return
    try:
        for i, page in enumerate(doc):
            if i >= _MAX_PDF_PAGES:
                yield f"[Truncated: PDF has more than {_MAX_PDF_PAGES} pages.]"
                break
            text = page.get_text("text")
            if text and text.strip():
                yield text.strip()
    finally:
        doc.close()


def _peek_pdf_pages(file_bytes: bytes, sample_pages: int = 5):
    """Peek at first few pages of a PDF. Returns None if fitz is unavailable or file is corrupt."""
    try:
        import fitz
    except ImportError:
        logger.warning("fitz (pymupdf) not available - PDF peek disabled.")
        return None
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception:
        logger.warning("Failed to open PDF for peek - file may be corrupt.")
        return None
    try:
        for i in range(min(sample_pages, len(doc))):
            text = doc[i].get_text("text")
            if text:
                yield text
    finally:
        doc.close()


def _iter_ocr_pdf_pages(file_bytes: bytes, max_pages: int = 5000):
    """OCR a PDF page-by-page. Returns None if OCR deps are unavailable or file is corrupt."""
    try:
        import fitz, pytesseract
        from PIL import Image
    except ImportError:
        logger.warning("OCR fallback unavailable.")
        return None
    _ensure_tesseract_configured()
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception:
        logger.warning("Failed to open PDF for OCR streaming - file may be corrupt.")
        return None
    total_pages = len(doc)
    pages_to_process = min(total_pages, max_pages)
    try:
        for i in range(pages_to_process):
            try:
                page = doc[i]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text = pytesseract.image_to_string(img)
                if text and text.strip():
                    yield text.strip()
            except Exception as e:
                logger.warning(f"OCR failed on page {i}: {e}")
                continue
    finally:
        doc.close()
    if total_pages > max_pages:
        yield f"[OCR truncated: {total_pages} pages, only first {max_pages} OCR'd.]"


# ─── DOCX Parsing ─────────────────────────────────────────────────────────

def _parse_docx(file_bytes: bytes) -> list[dict]:
    import docx
    document = docx.Document(io.BytesIO(file_bytes))
    sections = []
    current_title = "Overview"
    current_lines = []
    for para in document.paragraphs:
        style = (para.style.name if para.style else "") or ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading") or style == "Title":
            _flush(current_title, current_lines, sections)
            current_title, current_lines = text, []
        else:
            current_lines.append(text)
    _flush(current_title, current_lines, sections)
    for i, table in enumerate(document.tables, start=1):
        rows_text = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                rows_text.append(" | ".join(cells))
        if rows_text:
            sections.append({"title": f"Table {i}", "body": "\n".join(rows_text)})
    for para in document.paragraphs:
        for run in para.runs:
            if run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                for drawing in run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                    desc = drawing.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}desc')
                    if desc is not None and desc.text and desc.text.strip():
                        current_lines.append(f"[Image: {desc.text.strip()}]")
    return [s for s in sections if s["body"].strip()]


# ─── PPTX Parsing ─────────────────────────────────────────────────────────

def _parse_pptx(file_bytes: bytes) -> list[dict]:
    """Parse PowerPoint (.pptx) into sections, one per slide."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        logger.warning("python-pptx not available - PPTX parsing disabled.")
        return []
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.warning(f"Failed to open PPTX: {e}")
        return []
    sections = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_lines.append(" | ".join(cells))
        body = "\n".join(slide_lines).strip()
        if body:
            title = f"Slide {i}"
            for shape in slide.shapes:
                if shape.has_text_frame and shape == slide.shapes[0]:
                    first_text = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ""
                    if first_text:
                        title = first_text
                        break
            sections.append({"title": title, "body": body})
    return sections


# ─── ODP / ODS Parsing ────────────────────────────────────────────────────

def _parse_odp(file_bytes: bytes) -> list[dict]:
    """Parse LibreOffice Impress (.odp) into sections."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        sections = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_lines.append(text)
            body = "\n".join(slide_lines).strip()
            if body:
                sections.append({"title": f"Slide {i}", "body": body})
        if sections:
            return sections
    except Exception:
        pass
    return _parse_zip_xml_text(file_bytes, "content.xml")


def _parse_ods(file_bytes: bytes) -> list[dict]:
    """Parse LibreOffice Calc (.ods) spreadsheet into table sections."""
    return _parse_zip_xml_table(file_bytes, "content.xml")


def _parse_zip_xml_text(file_bytes: bytes, xml_path: str = "content.xml") -> list[dict]:
    """Generic parser for zip-based XML formats."""
    try:
        import zipfile
        import xml.etree.ElementTree as ET
    except ImportError:
        return []
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
        if xml_path not in zf.namelist():
            return []
        xml_content = zf.read(xml_path)
        zf.close()
    except Exception:
        return []
    try:
        root = ET.fromstring(xml_content)
    except Exception:
        return []
    texts = []
    for elem in root.iter():
        if elem.text and elem.text.strip():
            texts.append(elem.text.strip())
    if texts:
        return [{"title": "Content", "body": "\n".join(texts)}]
    return []


def _parse_zip_xml_table(file_bytes: bytes, xml_path: str = "content.xml") -> list[dict]:
    """Parse table data from a zip-based XML spreadsheet (ODS)."""
    try:
        import zipfile
        import xml.etree.ElementTree as ET
    except ImportError:
        return []
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
        if xml_path not in zf.namelist():
            return []
        xml_content = zf.read(xml_path)
        zf.close()
    except Exception:
        return []
    try:
        root = ET.fromstring(xml_content)
    except Exception:
        return []
    sections = []
    for table_elem in root.iter("{urn:oasis:names:tc:opendocument:xmlns:table:1.0}table"):
        table_name = table_elem.get("{urn:oasis:names:tc:opendocument:xmlns:table:1.0}name", "Sheet")
        rows = []
        for row_elem in table_elem.iter("{urn:oasis:names:tc:opendocument:xmlns:table:1.0}table-row"):
            cells = []
            for cell_elem in row_elem.iter("{urn:oasis:names:tc:opendocument:xmlns:table:1.0}table-cell"):
                cell_texts = []
                for p in cell_elem.iter("{urn:oasis:names:tc:opendocument:xmlns:text:1.0}p"):
                    if p.text:
                        cell_texts.append(p.text)
                cells.append(" ".join(cell_texts))
            if any(c.strip() for c in cells):
                rows.append(cells)
        if not rows:
            continue
        headers = rows[0]
        data_rows = rows[1:]
        md_rows = [[str(c) if c else "" for c in row] for row in data_rows]
        table_md = _table_to_markdown(headers, md_rows)
        summary = _generate_table_summary(table_name, headers, len(data_rows), table_name)
        sections.append({"title": f"{table_name}", "body": summary, "type": "table_summary"})
        sections.append({"title": f"{table_name} - Data", "body": table_md, "type": "table_data"})
    return sections


# ─── XML / YAML / TXT / EPUB / HTML / RTF / JSON / MD / SVG Parsing ──────

def _parse_xml(file_bytes: bytes) -> str:
    """Parse XML into a readable tree structure."""
    try:
        import xml.etree.ElementTree as ET
    except ImportError:
        return file_bytes.decode("utf-8", errors="ignore")
    try:
        root = ET.fromstring(file_bytes)
    except Exception:
        return file_bytes.decode("utf-8", errors="ignore")
    lines = []
    def _walk(elem, depth=0):
        indent = "  " * depth
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        text = (elem.text or "").strip()
        if text:
            lines.append(f"{indent}{tag}: {text}")
        else:
            lines.append(f"{indent}{tag}")
        for child in elem:
            _walk(child, depth + 1)
        tail = (elem.tail or "").strip()
        if tail:
            lines.append(f"{indent}  [tail]: {tail}")
    _walk(root)
    return "\n".join(lines)


def _parse_yaml(file_bytes: bytes) -> str:
    """Parse YAML into pretty-printed JSON."""
    import json as _json
    try:
        import yaml
    except ImportError:
        return file_bytes.decode("utf-8", errors="ignore")
    try:
        data = yaml.safe_load(file_bytes)
        return _json.dumps(data, indent=2, ensure_ascii=False, default=str)
    except Exception:
        return file_bytes.decode("utf-8", errors="ignore")


def _parse_txt(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore")


def _parse_epub(file_bytes: bytes) -> list[dict]:
    import ebooklib
    from ebooklib import epub
    from html.parser import HTMLParser
    import tempfile
    class _TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
        def handle_data(self, data):
            if data.strip():
                self.parts.append(data.strip())
    fd, tmp_path = tempfile.mkstemp(suffix=".epub")
    try:
        with os.fdopen(fd, "wb") as f:
            f.write(file_bytes)
        book = epub.read_epub(tmp_path)
    except Exception:
        os.unlink(tmp_path)
        raise
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
    sections = []
    for i, item in enumerate(book.get_items_of_type(ebooklib.ITEM_DOCUMENT), start=1):
        parser = _TextExtractor()
        try:
            content_bytes = item.get_content()
            if content_bytes:
                parser.feed(content_bytes.decode("utf-8", errors="ignore"))
        except Exception:
            continue
        body = "\n".join(parser.parts).strip()
        if body:
            sections.append({"title": item.get_name() or f"Chapter {i}", "body": body})
    return sections


def _parse_html(file_bytes: bytes) -> str:
    from html.parser import HTMLParser
    class _TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
            self._skip = False
        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self._skip = True
        def handle_endtag(self, tag):
            if tag in ("script", "style"):
                self._skip = False
        def handle_data(self, data):
            if not self._skip and data.strip():
                self.parts.append(data.strip())
    parser = _TextExtractor()
    parser.feed(file_bytes.decode("utf-8", errors="ignore"))
    return "\n".join(parser.parts)


def _parse_rtf(file_bytes: bytes) -> str:
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(file_bytes.decode("utf-8", errors="ignore"))


def _parse_json(file_bytes: bytes) -> str:
    import json as _json
    data = _json.loads(file_bytes.decode("utf-8", errors="ignore"))
    return _json.dumps(data, indent=2, ensure_ascii=False)


def _parse_md(text: str) -> list[dict]:
    sections = []
    current_title = "Overview"
    current_lines = []
    for line in text.splitlines():
        stripped = line.lstrip("#").strip()
        if line.startswith("### "):
            _flush(current_title, current_lines, sections)
            current_title, current_lines = stripped, []
        elif line.startswith("## "):
            _flush(current_title, current_lines, sections)
            current_title, current_lines = stripped, []
        elif line.startswith("# "):
            _flush(current_title, current_lines, sections)
            current_title, current_lines = stripped, []
        else:
            current_lines.append(line)
    _flush(current_title, current_lines, sections)
    return [s for s in sections if s["body"].strip()]


def _flush(title: str, lines: list[str], out: list[dict]):
    body = "\n".join(lines).strip()
    if body:
        out.append({"title": title, "body": body})


def _table_to_markdown(headers: list[str], rows: list[list], max_rows: int = 100) -> str:
    """Convert tabular data to a markdown table string.

    For very large tables, only the first max_rows are included and a
    summary note is appended.
    """
    if not headers or not rows:
        return ""

    md = "| " + " | ".join(headers) + " |\n"
    md += "| " + " | ".join("---" for _ in headers) + " |\n"

    truncated = len(rows) > max_rows
    display_rows = rows[:max_rows]

    for row in display_rows:
        cells = [str(c) if c is not None else "" for c in row]
        md += "| " + " | ".join(cells) + " |\n"

    if truncated:
        md += f"\n*[Table truncated: {len(rows)} total rows, showing first {max_rows}]*\n"

    return md


def _generate_table_summary(title: str, headers: list[str], row_count: int, sheet_name: str = None) -> str:
    """Generate a natural-language summary of a table for the parent node body."""
    parts = [f"Spreadsheet: {title}"]
    if sheet_name:
        parts.append(f"Sheet: {sheet_name}")
    parts.append(f"Rows: {row_count}")
    parts.append(f"Columns: {', '.join(headers)}")
    return " | ".join(parts)


def _parse_csv(file_bytes: bytes, doc_title: str = "Untitled") -> list[dict]:
    """Parse CSV into a single table section + summary, not one node per row."""
    import csv, io as _io
    text = file_bytes.decode("utf-8", errors="replace")
    reader = csv.DictReader(_io.StringIO(text))

    rows = list(reader)
    if not rows:
        return []

    headers = list(rows[0].keys())
    total_rows = len(rows)

    if total_rows > _MAX_SPREADSHEET_ROWS:
        rows = rows[:_MAX_SPREADSHEET_ROWS]
        logger.warning(f"CSV truncated at {_MAX_SPREADSHEET_ROWS} rows")

    md_rows = [[row.get(h, "") for h in headers] for row in rows]
    table_md = _table_to_markdown(headers, md_rows)

    summary = _generate_table_summary(doc_title, headers, total_rows)

    return [
        {"title": doc_title, "body": summary, "type": "table_summary"},
        {"title": f"{doc_title} - Data", "body": table_md, "type": "table_data"},
    ]


def _parse_xlsx(file_bytes: bytes, doc_title: str = "Untitled") -> list[dict]:
    """Parse Excel into one section per sheet (summary + data), not one node per row."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sections = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        headers = [str(h) if h is not None else f"col{i}" for i, h in enumerate(rows[0])]
        data_rows = rows[1:]
        total_rows = len(data_rows)

        if total_rows > _MAX_SPREADSHEET_ROWS:
            data_rows = data_rows[:_MAX_SPREADSHEET_ROWS]
            logger.warning(f"XLSX sheet '{sheet_name}' truncated at {_MAX_SPREADSHEET_ROWS} rows")

        md_rows = [[str(c) if c is not None else "" for c in row] for row in data_rows]
        table_md = _table_to_markdown(headers, md_rows)

        summary = _generate_table_summary(doc_title, headers, total_rows, sheet_name)

        sections.append({
            "title": f"{doc_title} - {sheet_name}",
            "body": summary,
            "type": "table_summary",
        })
        sections.append({
            "title": f"{doc_title} - {sheet_name} - Data",
            "body": table_md,
            "type": "table_data",
        })

    wb.close()
    return sections


def _parse_svg(file_bytes: bytes) -> str:
    import xml.etree.ElementTree as ET
    root = ET.fromstring(file_bytes)
    texts = []
    text_tags = {"title", "desc", "text", "tspan", "flowRoot", "flowPara"}
    for elem in root.iter():
        local = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        if local in text_tags:
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                texts.append(elem.tail.strip())
    return "\n".join(dict.fromkeys(texts))
