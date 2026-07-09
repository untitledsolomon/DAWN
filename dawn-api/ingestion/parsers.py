"""
DAWN Document Parsers — multi-format text extraction engine.

Supports:
  EPUB2 / EPUB3 (ebooklib + lxml XHTML parsing)
  MOBI / AZW3 / AZW (conversion via calibre or mobi-python)
  FB2 (FictionBook XML)
  DjVu (via djvulibre or pure-Python fallback)
  ODT (LibreOffice Writer — zip-based XML)
  LaTeX (.tex, .sty, .cls)
  Plain text (.txt, .log, .ini, .cfg, .conf)
  Markdown (.md, .markdown)
  HTML (.html, .htm)
  RTF (via striprtf)
  SVG (XML text extraction)
  XML / YAML / JSON
  CSV / Excel (.xlsx, .xls)
  Word (.docx)
  PowerPoint (.pptx)
  ODP / ODS (LibreOffice Impress / Calc)

All parsers return a list[dict] with {"title": str, "body": str} sections,
or a plain text string for formats that don't have natural section breaks.
"""

import io
import os
import re
import json
import logging
import zipfile
import tempfile
import xml.etree.ElementTree as ET
from html.parser import HTMLParser
from typing import Optional

logger = logging.getLogger(__name__)

_MAX_SPREADSHEET_ROWS = 100_000
_MAX_TABLE_BODY_CHARS = 50_000


def _flush(title: str, lines: list[str], out: list[dict]):
    body = "\n".join(lines).strip()
    if body:
        out.append({"title": title, "body": body})


def _table_to_markdown(headers: list[str], rows: list[list], max_rows: int = 100) -> str:
    if not headers or not rows:
        return ""
    md = "| " + " | ".join(headers) + " |\n"
    md += "| " + " | ".join("---" for _ in headers) + " |\n"
    truncated = len(rows) > max_rows
    display_rows = rows[:max_rows]
    for row in display_rows:
        cells = [str(c) if c is not None else "" for c in row]
        md += "| " + " | ".join(cells) + " |\n"
    if truncated:
        md += f"\n*[Table truncated: {len(rows)} total rows, showing first {max_rows}]*\n"
    return md


def _generate_table_summary(title: str, headers: list[str], row_count: int, sheet_name: str = None) -> str:
    parts = [f"Spreadsheet: {title}"]
    if sheet_name:
        parts.append(f"Sheet: {sheet_name}")
    parts.append(f"Rows: {row_count}")
    parts.append(f"Columns: {', '.join(headers)}")
    return " | ".join(parts)


# ═══════════════════════════════════════════════════════════════════════════
# EPUB2 / EPUB3 Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_epub(file_bytes: bytes) -> list[dict]:
    """
    Parse EPUB2 and EPUB3 files into sections.

    EPUB2 uses NCX for the table of contents.
    EPUB3 uses the nav document (EPUB3 Navigation Document).
    Both store content as XHTML/HTML files in the zip archive.

    This parser handles both by:
    1. Reading the OPF manifest to find all content documents
    2. Reading the NCX (EPUB2) or nav (EPUB3) for chapter titles
    3. Extracting text from each content document via HTML parsing
    """
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError:
        logger.warning("ebooklib not available — EPUB parsing disabled.")
        return _parse_epub_fallback(file_bytes)

    try:
        fd, tmp_path = tempfile.mkstemp(suffix=".epub")
        with os.fdopen(fd, "wb") as f:
            f.write(file_bytes)
        book = epub.read_epub(tmp_path)
    except Exception as e:
        logger.warning(f"Failed to open EPUB: {e}")
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
        return _parse_epub_fallback(file_bytes)
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

    sections = []

    # Build a map of item ID -> (href, media-type, content)
    item_map = {}
    for item in book.get_items():
        item_map[item.get_id()] = {
            "href": item.get_name(),
            "media_type": getattr(item, "get_media_type", lambda: "image/unknown")(),
            "content": item.get_content(),
        }

    # Get chapter ordering from the spine
    spine_order = []
    try:
        spine = book.get_spine()
        for item_id, linear in spine:
            if item_id in item_map:
                spine_order.append(item_id)
    except Exception:
        pass

    # Try to get chapter titles from NCX (EPUB2)
    ncx_titles = {}
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
        ncx_file = None
        for name in zf.namelist():
            if name.endswith(".ncx"):
                ncx_file = name
                break
        if ncx_file:
            ncx_content = zf.read(ncx_file).decode("utf-8", errors="ignore")
            ncx_titles = _parse_ncx(ncx_content)
        zf.close()
    except Exception:
        pass

    # Try EPUB3 nav document if NCX didn't yield titles
    if not ncx_titles:
        try:
            for item_id, info in item_map.items():
                href = info["href"] or ""
                if "nav" in href.lower() or "toc" in href.lower():
                    content = info["content"]
                    if content:
                        ncx_titles = _parse_epub3_nav(content.decode("utf-8", errors="ignore"))
                        break
        except Exception:
            pass

    # Extract content from each spine item
    if spine_order:
        for item_id in spine_order:
            info = item_map.get(item_id)
            if not info:
                continue
            href = info["href"] or f"chapter_{item_id}"
            content_bytes = info.get("content")
            if not content_bytes:
                continue
            text = _extract_html_text(content_bytes.decode("utf-8", errors="ignore"))
            if text.strip():
                title = ncx_titles.get(href, ncx_titles.get(item_id,
                    os.path.splitext(os.path.basename(href))[0]))
                sections.append({"title": title, "body": text.strip()})
    else:
        # Fallback: iterate all document items
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            content_bytes = item.get_content()
            if not content_bytes:
                continue
            text = _extract_html_text(content_bytes.decode("utf-8", errors="ignore"))
            if text.strip():
                name = item.get_name() or "chapter"
                title = ncx_titles.get(name, os.path.splitext(os.path.basename(name))[0])
                sections.append({"title": title, "body": text.strip()})

    return sections


def _parse_ncx(ncx_content: str) -> dict[str, str]:
    """Parse EPUB2 NCX table of contents. Returns {href: title}."""
    result = {}
    try:
        root = ET.fromstring(ncx_content)
        for nav_point in root.iter("{http://www.daisy.org/z3986/2005/ncx/}navPoint"):
            text_el = nav_point.find(".//{http://www.daisy.org/z3986/2005/ncx/}text")
            content_el = nav_point.find(".//{http://www.daisy.org/z3986/2005/ncx/}content")
            if text_el is not None and content_el is not None:
                title = text_el.text or ""
                src = content_el.get("src", "")
                if src and title:
                    href = src.split("#")[0]
                    result[href] = title
    except Exception as e:
        logger.debug(f"NCX parsing failed: {e}")
    return result


def _parse_epub3_nav(nav_content: str) -> dict[str, str]:
    """Parse EPUB3 Navigation Document. Returns {href: title}."""
    result = {}
    try:
        root = ET.fromstring(nav_content)
        for a in root.iter("{http://www.w3.org/1999/xhtml}a"):
            href = a.get("href", "")
            title = a.text or ""
            if href and title:
                href = href.split("#")[0]
                result[href] = title
        if not result:
            for a in root.iter("a"):
                href = a.get("href", "")
                title = a.text or ""
                if href and title:
                    href = href.split("#")[0]
                    result[href] = title
    except Exception as e:
        logger.debug(f"EPUB3 nav parsing failed: {e}")
    return result


def _extract_html_text(html_content: str) -> str:
    """Extract clean text from HTML/XHTML content."""
    class _Extractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
            self._skip = False
            self._skip_tags = {"script", "style", "nav", "header", "footer"}
        def handle_starttag(self, tag, attrs):
            if tag in self._skip_tags:
                self._skip = True
        def handle_endtag(self, tag):
            if tag in self._skip_tags:
                self._skip = False
        def handle_data(self, data):
            if not self._skip and data.strip():
                self.parts.append(data.strip())
    parser = _Extractor()
    try:
        parser.feed(html_content)
    except Exception:
        return html_content
    return "\n".join(parser.parts)


def _parse_epub_fallback(file_bytes: bytes) -> list[dict]:
    """
    Fallback EPUB parser that reads the zip directly.
    Works for both EPUB2 and EPUB3 without ebooklib.
    """
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
    except Exception:
        return []

    sections = []

    # Find the OPF file via META-INF/container.xml
    opf_path = None
    try:
        container_xml = zf.read("META-INF/container.xml").decode("utf-8", errors="ignore")
        root = ET.fromstring(container_xml)
        for rf in root.iter():
            if "rootfile" in rf.tag.lower() or rf.tag.endswith("rootfile"):
                opf_path = rf.get("full-path")
                break
    except Exception:
        pass

    if not opf_path:
        for name in zf.namelist():
            if name.endswith(".opf"):
                opf_path = name
                break

    if not opf_path:
        zf.close()
        return []

    try:
        opf_content = zf.read(opf_path).decode("utf-8", errors="ignore")
        opf_root = ET.fromstring(opf_content)
    except Exception:
        zf.close()
        return []

    base_dir = os.path.dirname(opf_path)

    # Build id -> href map from manifest
    id_to_href = {}
    content_files = set()
    for item in opf_root.iter():
        tag = item.tag.split("}")[-1] if "}" in item.tag else item.tag
        if tag == "item":
            media_type = item.get("media-type", "")
            href = item.get("href", "")
            item_id = item.get("id")
            if item_id:
                id_to_href[item_id] = href
            if "html" in media_type.lower() or "xhtml" in media_type.lower():
                full_path = os.path.normpath(os.path.join(base_dir, href))
                content_files.add(full_path)

    # Get spine ordering
    spine_items = []
    for spine in opf_root.iter():
        tag = spine.tag.split("}")[-1] if "}" in spine.tag else spine.tag
        if tag == "spine":
            for itemref in spine:
                iref_tag = itemref.tag.split("}")[-1] if "}" in itemref.tag else itemref.tag
                if iref_tag == "itemref":
                    idref = itemref.get("idref")
                    if idref:
                        spine_items.append(idref)

    ordered_files = []
    seen = set()
    for idref in spine_items:
        href = id_to_href.get(idref)
        if href:
            full_path = os.path.normpath(os.path.join(base_dir, href))
            if full_path not in seen:
                ordered_files.append(full_path)
                seen.add(full_path)

    for f in sorted(content_files):
        if f not in seen:
            ordered_files.append(f)
            seen.add(f)

    for fpath in ordered_files:
        try:
            content = zf.read(fpath).decode("utf-8", errors="ignore")
        except Exception:
            continue
        text = _extract_html_text(content)
        if text.strip():
            title = os.path.splitext(os.path.basename(fpath))[0]
            sections.append({"title": title, "body": text.strip()})

    zf.close()
    return sections


# ═══════════════════════════════════════════════════════════════════════════
# MOBI / AZW3 / AZW Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_mobi(file_bytes: bytes) -> list[dict]:
    """
    Parse MOBI / AZW3 / AZW files.

    Strategy (tried in order):
    1. Convert via calibre's ebook-convert CLI (if available)
    2. Use mobi-python (mobi package) for header extraction
    3. Fallback: extract raw text heuristically
    """
    # Strategy 1: calibre conversion
    try:
        return _parse_mobi_calibre(file_bytes)
    except Exception as e:
        logger.debug(f"MOBI calibre conversion failed: {e}")

    # Strategy 2: mobi-python
    try:
        return _parse_mobi_python(file_bytes)
    except Exception as e:
        logger.debug(f"MOBI mobi-python parsing failed: {e}")

    # Strategy 3: raw text heuristic
    return _parse_mobi_heuristic(file_bytes)


def _parse_mobi_calibre(file_bytes: bytes) -> list[dict]:
    """Convert MOBI to text via calibre's ebook-convert."""
    import subprocess
    fd_in, tmp_in = tempfile.mkstemp(suffix=".mobi")
    fd_out, tmp_out = tempfile.mkstemp(suffix=".txt")
    os.close(fd_out)
    try:
        with os.fdopen(fd_in, "wb") as f:
            f.write(file_bytes)
        result = subprocess.run(
            ["ebook-convert", tmp_in, tmp_out],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode != 0:
            raise RuntimeError(f"ebook-convert failed: {result.stderr[:200]}")
        with open(tmp_out, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        if text.strip():
            return [{"title": "Content", "body": text.strip()}]
    finally:
        try:
            os.unlink(tmp_in)
        except Exception:
            pass
        try:
            os.unlink(tmp_out)
        except Exception:
            pass
    return []


def _parse_mobi_python(file_bytes: bytes) -> list[dict]:
    """Parse MOBI using the mobi Python package."""
    import struct
    # mobi package: pip install mobi
    # It returns (filepath, status) after extracting
    try:
        import mobi
        fd, tmp_path = tempfile.mkstemp(suffix=".mobi")
        with os.fdopen(fd, "wb") as f:
            f.write(file_bytes)
        tempdir, status = mobi.extract(tmp_path)
        if status is None:
            # mobi.extract returns (tempdir, None) on success
            sections = []
            for root, dirs, files in os.walk(tempdir):
                for fname in sorted(files):
                    fpath = os.path.join(root, fname)
                    if fname.endswith(".html") or fname.endswith(".xhtml"):
                        with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                            text = _extract_html_text(f.read())
                        if text.strip():
                            title = os.path.splitext(fname)[0]
                            sections.append({"title": title, "body": text.strip()})
                    elif fname.endswith(".txt"):
                        with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                            text = f.read()
                        if text.strip():
                            sections.append({"title": fname, "body": text.strip()})
            import shutil
            shutil.rmtree(tempdir, ignore_errors=True)
            if sections:
                return sections
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
    except ImportError:
        logger.debug("mobi package not available")
    except Exception as e:
        logger.debug(f"mobi-python extraction failed: {e}")
    return []


def _parse_mobi_heuristic(file_bytes: bytes) -> list[dict]:
    """
    Heuristic MOBI text extraction.

    MOBI files have a PDB header (first 78 bytes), then a MOBI header.
    Text records follow the headers. This extracts concatenated text records.
    """
    try:
        if len(file_bytes) < 78:
            return []
        # Check for MOBI magic in the header
        mobi_header_start = struct.unpack(">I", file_bytes[76:80])[0] if len(file_bytes) > 80 else 0
        if mobi_header_start == 0:
            return []

        # Try to find text by looking for readable ASCII/UTF-8 sequences
        # after the MOBI header
        text_start = mobi_header_start + 16  # MOBI header is at least 16 bytes
        if text_start >= len(file_bytes):
            return []

        raw = file_bytes[text_start:]
        # Decode with error-tolerant approach
        text = raw.decode("utf-8", errors="replace")
        # Clean up: remove null bytes and control characters (keep newlines)
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
        text = re.sub(r"\s+", " ", text).strip()
        if len(text) > 100:
            return [{"title": "Content", "body": text}]
    except Exception as e:
        logger.debug(f"MOBI heuristic failed: {e}")
    return []


# ═══════════════════════════════════════════════════════════════════════════
# FB2 (FictionBook) Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_fb2(file_bytes: bytes) -> list[dict]:
    """
    Parse FictionBook (.fb2) XML format.

    FB2 is an XML format with:
    - <description> (metadata)
    - <body> (main content with <section> elements)
    - <binary> (embedded images, skipped)
    """
    try:
        root = ET.fromstring(file_bytes)
    except Exception as e:
        logger.warning(f"FB2 XML parsing failed: {e}")
        return []

    # Handle FB2 namespace
    ns = ""
    if root.tag.startswith("{"):
        ns = root.tag[:root.tag.index("}") + 1]

    sections = []

    # Extract title info from description
    book_title = "FB2 Document"
    try:
        title_info = root.find(f".//{ns}title-info")
        if title_info is not None:
            bt = title_info.find(f"{ns}book-title")
            if bt is not None and bt.text:
                book_title = bt.text
    except Exception:
        pass

    # Extract sections from body
    body = root.find(f"{ns}body")
    if body is None:
        return []

    for section in body.findall(f".//{ns}section"):
        section_title = ""
        section_parts = []

        for child in section:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "title":
                # Collect all text in the title element
                title_texts = []
                for t in child.iter():
                    if t.text and t.text.strip():
                        title_texts.append(t.text.strip())
                section_title = " ".join(title_texts)
            elif tag in ("p", "poem", "epigraph", "cite", "subtitle"):
                text = _get_all_text(child)
                if text.strip():
                    section_parts.append(text.strip())

        body_text = "\n\n".join(section_parts)
        if body_text.strip():
            sections.append({
                "title": section_title or f"Section {len(sections) + 1}",
                "body": body_text,
            })

    if not sections:
        # Fallback: extract all text from body
        all_text = _get_all_text(body)
        if all_text.strip():
            sections.append({"title": book_title, "body": all_text.strip()})

    return sections


def _get_all_text(element) -> str:
    """Get all text content from an XML element recursively."""
    texts = []
    if element.text and element.text.strip():
        texts.append(element.text.strip())
    for child in element:
        texts.append(_get_all_text(child))
        if child.tail and child.tail.strip():
            texts.append(child.tail.strip())
    return " ".join(texts)


# ═══════════════════════════════════════════════════════════════════════════
# DjVu Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_djvu(file_bytes: bytes) -> list[dict]:
    """
    Parse DjVu files.

    Strategy:
    1. Use djvutxt CLI from djvulibre (if available)
    2. Fallback: extract text from the DjVu bundled XML/annotations
    """
    try:
        return _parse_djvu_cli(file_bytes)
    except Exception as e:
        logger.debug(f"DjVu CLI parsing failed: {e}")

    try:
        return _parse_djvu_heuristic(file_bytes)
    except Exception as e:
        logger.debug(f"DjVu heuristic parsing failed: {e}")

    return []


def _parse_djvu_cli(file_bytes: bytes) -> list[dict]:
    """Extract text from DjVu using djvutxt."""
    import subprocess
    fd, tmp_path = tempfile.mkstemp(suffix=".djvu")
    try:
        with os.fdopen(fd, "wb") as f:
            f.write(file_bytes)
        result = subprocess.run(
            ["djvutxt", tmp_path],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return [{"title": "Content", "body": result.stdout.strip()}]
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
    return []


def _parse_djvu_heuristic(file_bytes: bytes) -> list[dict]:
    """
    Heuristic DjVu text extraction.

    DjVu files contain chunks with 4-byte IDs (like 'FORM', 'DJVI', 'INFO',
    'TXTz', etc.). Text is stored in 'TXTz' chunks (zlib-compressed) or
    'TXTa' chunks (plain text annotations).
    """
    try:
        import zlib
    except ImportError:
        return []

    texts = []
    pos = 0
    while pos + 8 < len(file_bytes):
        chunk_id = file_bytes[pos:pos+4].decode("ascii", errors="ignore")
        chunk_len = struct.unpack(">I", file_bytes[pos+4:pos+8])[0]
        chunk_data = file_bytes[pos+8:pos+8+chunk_len]

        if chunk_id == "TXTz":
            try:
                decompressed = zlib.decompress(chunk_data)
                text = decompressed.decode("utf-8", errors="replace")
                texts.append(text)
            except Exception:
                pass
        elif chunk_id == "TXTa":
            try:
                text = chunk_data.decode("utf-8", errors="replace")
                texts.append(text)
            except Exception:
                pass

        pos += 8 + chunk_len
        if chunk_len % 2 != 0:
            pos += 1  # padding

    if texts:
        combined = "\n".join(texts)
        combined = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", combined)
        combined = re.sub(r"\s+", " ", combined).strip()
        if len(combined) > 50:
            return [{"title": "Content", "body": combined}]
    return []


# ═══════════════════════════════════════════════════════════════════════════
# ODT (LibreOffice Writer) Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_odt(file_bytes: bytes) -> list[dict]:
    """
    Parse ODT (OpenDocument Text) files.

    ODT is a zip containing content.xml with the document body.
    """
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
    except Exception:
        return []

    if "content.xml" not in zf.namelist():
        zf.close()
        return []

    try:
        xml_content = zf.read("content.xml")
        zf.close()
    except Exception:
        zf.close()
        return []

    try:
        root = ET.fromstring(xml_content)
    except Exception:
        return []

    # ODF namespaces
    ns = {
        "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
        "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
        "style": "urn:oasis:names:tc:opendocument:xmlns:style:1.0",
    }

    sections = []
    current_title = "Content"
    current_lines = []

    body = root.find(".//office:body", ns)
    if body is None:
        body = root.find(".//{urn:oasis:names:tc:opendocument:xmlns:office:1.0}body")

    if body is None:
        return []

    text_elem = body.find(".//office:text", ns)
    if text_elem is None:
        text_elem = body.find(".//{urn:oasis:names:tc:opendocument:xmlns:office:1.0}text")

    if text_elem is None:
        return []

    for child in text_elem:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "h":
            # Heading
            _flush(current_title, current_lines, sections)
            current_title = _get_all_text(child) or f"Heading {len(sections) + 1}"
            current_lines = []
        elif tag == "p":
            text = _get_all_text(child)
            if text.strip():
                current_lines.append(text.strip())
        elif tag == "list":
            for list_item in child.findall(".//text:list-item", ns):
                text = _get_all_text(list_item)
                if text.strip():
                    current_lines.append(f"  - {text.strip()}")
        elif tag == "table":
            # Extract table as markdown
            rows = []
            for table_row in child.findall(".//text:table-row", ns):
                cells = []
                for cell in table_row.findall(".//text:table-cell", ns):
                    cell_text = _get_all_text(cell)
                    cells.append(cell_text)
                if any(c.strip() for c in cells):
                    rows.append(cells)
            if rows:
                headers = rows[0]
                md_rows = rows[1:]
                current_lines.append("\n" + _table_to_markdown(headers, md_rows))

    _flush(current_title, current_lines, sections)
    return [s for s in sections if s["body"].strip()]


# ═══════════════════════════════════════════════════════════════════════════
# LaTeX Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_latex(file_bytes: bytes) -> list[dict]:
    """
    Parse LaTeX (.tex, .sty, .cls) files into sections.

    Handles:
    - \section, \subsection, \subsubsection
    - \chapter, \part
    - \begin{abstract} ... \end{abstract}
    - Basic math mode removal
    - \cite, \ref, \label stripping
    """
    text = file_bytes.decode("utf-8", errors="ignore")

    # Remove LaTeX comments (but not \% escaped ones)
    text = re.sub(r"(?<!\\)%.*$", "", text, flags=re.MULTILINE)

    sections = []
    current_title = "Preamble"
    current_lines = []

    for line in text.splitlines():
        stripped = line.strip()

        # Skip empty lines and control sequences
        if not stripped:
            continue

        # Detect section commands
        section_match = re.match(
            r"\\(chapter|part|section|subsection|subsubsection)\*?\s*\{([^}]*)\}",
            stripped,
        )
        if section_match:
            _flush_latex(current_title, current_lines, sections)
            current_title = section_match.group(2)
            current_lines = []
            continue

        # Detect \begin{abstract}
        if stripped.startswith(r"\begin{abstract}"):
            _flush_latex(current_title, current_lines, sections)
            current_title = "Abstract"
            current_lines = []
            continue

        if stripped.startswith(r"\end{abstract}"):
            _flush_latex(current_title, current_lines, sections)
            current_title = "Content"
            current_lines = []
            continue

        # Skip other begin/end blocks (we don't want raw LaTeX)
        if stripped.startswith(r"\begin{") or stripped.startswith(r"\end{"):
            continue

        # Skip control sequences at line start
        if stripped.startswith("\\") and not stripped.startswith("\\%"):
            # It's a command — skip unless it's text
            if not any(stripped.startswith(cmd) for cmd in [r"\textit", r"\textbf", r"\emph", r"\texttt"]):
                continue

        # Clean the line
        cleaned = _clean_latex_line(stripped)
        if cleaned:
            current_lines.append(cleaned)

    _flush_latex(current_title, current_lines, sections)
    return [s for s in sections if s["body"].strip()]


def _clean_latex_line(line: str) -> str:
    """Clean a single line of LaTeX text."""
    # Remove \label{...}
    line = re.sub(r"\\label\{[^}]*\}", "", line)
    # Remove \ref{...} and \cite{...}
    line = re.sub(r"\\(ref|cite|pageref)\{[^}]*\}", "", line)
    # Remove \index{...}
    line = re.sub(r"\\index\{[^}]*\}", "", line)
    # Remove \footnote{...}
    line = re.sub(r"\\footnote\{[^}]*\}", "", line)
    # Convert \textit{...}, \textbf{...}, \emph{...} to just text
    line = re.sub(r"\\(textit|textbf|emph|texttt|textsc)\{([^}]*)\}", r"\2", line)
    # Remove \url{...}
    line = re.sub(r"\\url\{[^}]*\}", "", line)
    # Remove \includegraphics[...]{...}
    line = re.sub(r"\\includegraphics(?:\[[^\]]*\])?\{[^}]*\}", "", line)
    # Remove math mode: $...$ and $$...$$
    line = re.sub(r"\$\$[^$]*\$\$", "", line)
    line = re.sub(r"\$[^$]*\$", "", line)
    # Remove \[ ... \] display math
    line = re.sub(r"\\\[.*?\\\]", "", line)
    # Collapse multiple spaces
    line = re.sub(r"\s+", " ", line).strip()
    return line


def _flush_latex(title: str, lines: list[str], out: list[dict]):
    body = "\n".join(lines).strip()
    if body:
        out.append({"title": title, "body": body})


# ═══════════════════════════════════════════════════════════════════════════
# Plain Text Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_text(file_bytes: bytes) -> list[dict]:
    """Parse plain text files. Returns a single section."""
    text = file_bytes.decode("utf-8", errors="ignore").strip()
    if not text:
        return []
    return [{"title": "Content", "body": text}]


# ═══════════════════════════════════════════════════════════════════════════
# Markdown Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_markdown(text: str) -> list[dict]:
    """Parse Markdown into sections by heading."""
    sections = []
    current_title = "Overview"
    current_lines = []
    for line in text.splitlines():
        stripped = line.lstrip("#").strip()
        if line.startswith("### "):
            _flush(current_title, current_lines, sections)
            current_title, current_lines = stripped, []
        elif line.startswith("## "):
            _flush(current_title, current_lines, sections)
            current_title, current_lines = stripped, []
        elif line.startswith("# "):
            _flush(current_title, current_lines, sections)
            current_title, current_lines = stripped, []
        else:
            current_lines.append(line)
    _flush(current_title, current_lines, sections)
    return [s for s in sections if s["body"].strip()]


# ═══════════════════════════════════════════════════════════════════════════
# HTML Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_html(file_bytes: bytes) -> list[dict]:
    """Parse HTML into clean text."""
    text = _extract_html_text(file_bytes.decode("utf-8", errors="ignore"))
    if text.strip():
        return [{"title": "Content", "body": text.strip()}]
    return []


# ═══════════════════════════════════════════════════════════════════════════
# RTF Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_rtf(file_bytes: bytes) -> list[dict]:
    """Parse RTF into clean text."""
    try:
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(file_bytes.decode("utf-8", errors="ignore"))
        if text.strip():
            return [{"title": "Content", "body": text.strip()}]
    except Exception as e:
        logger.warning(f"RTF parsing failed: {e}")
    return []


# ═══════════════════════════════════════════════════════════════════════════
# SVG Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_svg(file_bytes: bytes) -> list[dict]:
    """Extract text from SVG files."""
    try:
        root = ET.fromstring(file_bytes)
    except Exception:
        return []
    texts = []
    text_tags = {"title", "desc", "text", "tspan", "flowRoot", "flowPara"}
    for elem in root.iter():
        local = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        if local in text_tags:
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                texts.append(elem.tail.strip())
    combined = "\n".join(dict.fromkeys(texts))
    if combined.strip():
        return [{"title": "Content", "body": combined.strip()}]
    return []


# ═══════════════════════════════════════════════════════════════════════════
# XML / YAML / JSON Parsers
# ═══════════════════════════════════════════════════════════════════════════

def parse_xml(file_bytes: bytes) -> list[dict]:
    """Parse XML into a readable tree structure."""
    try:
        root = ET.fromstring(file_bytes)
    except Exception:
        text = file_bytes.decode("utf-8", errors="ignore")
        return [{"title": "Content", "body": text}] if text.strip() else []
    lines = []
    def _walk(elem, depth=0):
        indent = "  " * depth
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        text = (elem.text or "").strip()
        if text:
            lines.append(f"{indent}{tag}: {text}")
        else:
            lines.append(f"{indent}{tag}")
        for child in elem:
            _walk(child, depth + 1)
        tail = (elem.tail or "").strip()
        if tail:
            lines.append(f"{indent}  [tail]: {tail}")
    _walk(root)
    return [{"title": "Content", "body": "\n".join(lines)}]


def parse_yaml(file_bytes: bytes) -> list[dict]:
    """Parse YAML into pretty-printed JSON."""
    try:
        import yaml
        data = yaml.safe_load(file_bytes)
        text = json.dumps(data, indent=2, ensure_ascii=False, default=str)
        return [{"title": "Content", "body": text}]
    except ImportError:
        text = file_bytes.decode("utf-8", errors="ignore")
        return [{"title": "Content", "body": text}] if text.strip() else []


def parse_json(file_bytes: bytes) -> list[dict]:
    """Parse JSON into pretty-printed text."""
    try:
        data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
        text = json.dumps(data, indent=2, ensure_ascii=False)
        return [{"title": "Content", "body": text}]
    except Exception:
        text = file_bytes.decode("utf-8", errors="ignore")
        return [{"title": "Content", "body": text}] if text.strip() else []


# ═══════════════════════════════════════════════════════════════════════════
# CSV Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_csv(file_bytes: bytes, doc_title: str = "Untitled") -> list[dict]:
    """Parse CSV into table sections (summary + data)."""
    import csv as _csv
    import io as _io
    text = file_bytes.decode("utf-8", errors="replace")
    reader = _csv.DictReader(_io.StringIO(text))
    rows = list(reader)
    if not rows:
        return []
    headers = list(rows[0].keys())
    total_rows = len(rows)
    if total_rows > _MAX_SPREADSHEET_ROWS:
        rows = rows[:_MAX_SPREADSHEET_ROWS]
        logger.warning(f"CSV truncated at {_MAX_SPREADSHEET_ROWS} rows")
    md_rows = [[row.get(h, "") for h in headers] for row in rows]
    table_md = _table_to_markdown(headers, md_rows)
    summary = _generate_table_summary(doc_title, headers, total_rows)
    return [
        {"title": doc_title, "body": summary, "type": "table_summary"},
        {"title": f"{doc_title} - Data", "body": table_md, "type": "table_data"},
    ]


# ═══════════════════════════════════════════════════════════════════════════
# Excel Parser (.xlsx, .xls)
# ═══════════════════════════════════════════════════════════════════════════

def parse_xlsx(file_bytes: bytes, doc_title: str = "Untitled") -> list[dict]:
    """Parse Excel into one section per sheet (summary + data)."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sections = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        headers = [str(h) if h is not None else f"col{i}" for i, h in enumerate(rows[0])]
        data_rows = rows[1:]
        total_rows = len(data_rows)
        if total_rows > _MAX_SPREADSHEET_ROWS:
            data_rows = data_rows[:_MAX_SPREADSHEET_ROWS]
            logger.warning(f"XLSX sheet '{sheet_name}' truncated at {_MAX_SPREADSHEET_ROWS} rows")
        md_rows = [[str(c) if c is not None else "" for c in row] for row in data_rows]
        table_md = _table_to_markdown(headers, md_rows)
        summary = _generate_table_summary(doc_title, headers, total_rows, sheet_name)
        sections.append({
            "title": f"{doc_title} - {sheet_name}",
            "body": summary,
            "type": "table_summary",
        })
        sections.append({
            "title": f"{doc_title} - {sheet_name} - Data",
            "body": table_md,
            "type": "table_data",
        })
    wb.close()
    return sections


# ═══════════════════════════════════════════════════════════════════════════
# Word Parser (.docx)
# ═══════════════════════════════════════════════════════════════════════════

def parse_docx(file_bytes: bytes) -> list[dict]:
    """Parse Word (.docx) into sections by heading."""
    import docx
    document = docx.Document(io.BytesIO(file_bytes))
    sections = []
    current_title = "Overview"
    current_lines = []
    for para in document.paragraphs:
        style = (para.style.name if para.style else "") or ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading") or style == "Title":
            _flush(current_title, current_lines, sections)
            current_title, current_lines = text, []
        else:
            current_lines.append(text)
    _flush(current_title, current_lines, sections)
    for i, table in enumerate(document.tables, start=1):
        rows_text = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                rows_text.append(" | ".join(cells))
        if rows_text:
            sections.append({"title": f"Table {i}", "body": "\n".join(rows_text)})
    return [s for s in sections if s["body"].strip()]


# ═══════════════════════════════════════════════════════════════════════════
# PowerPoint Parser (.pptx)
# ═══════════════════════════════════════════════════════════════════════════

def parse_pptx(file_bytes: bytes) -> list[dict]:
    """Parse PowerPoint (.pptx) into sections, one per slide."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not available — PPTX parsing disabled.")
        return []
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.warning(f"Failed to open PPTX: {e}")
        return []
    sections = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_lines.append(" | ".join(cells))
        body = "\n".join(slide_lines).strip()
        if body:
            title = f"Slide {i}"
            for shape in slide.shapes:
                if shape.has_text_frame and shape == slide.shapes[0]:
                    first_text = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ""
                    if first_text:
                        title = first_text
                        break
            sections.append({"title": title, "body": body})
    return sections


# ═══════════════════════════════════════════════════════════════════════════
# ODP (LibreOffice Impress) Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_odp(file_bytes: bytes) -> list[dict]:
    """Parse LibreOffice Impress (.odp) into sections."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        sections = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_lines.append(text)
            body = "\n".join(slide_lines).strip()
            if body:
                sections.append({"title": f"Slide {i}", "body": body})
        if sections:
            return sections
    except Exception:
        pass
    return _parse_zip_xml_text(file_bytes, "content.xml")


# ═══════════════════════════════════════════════════════════════════════════
# ODS (LibreOffice Calc) Parser
# ═══════════════════════════════════════════════════════════════════════════

def parse_ods(file_bytes: bytes) -> list[dict]:
    """Parse LibreOffice Calc (.ods) spreadsheet into table sections."""
    return _parse_zip_xml_table(file_bytes, "content.xml")


# ═══════════════════════════════════════════════════════════════════════════
# Generic Zip-XML Helpers
# ═══════════════════════════════════════════════════════════════════════════

def _parse_zip_xml_text(file_bytes: bytes, xml_path: str = "content.xml") -> list[dict]:
    """Generic parser for zip-based XML formats (text extraction)."""
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
        if xml_path not in zf.namelist():
            zf.close()
            return []
        xml_content = zf.read(xml_path)
        zf.close()
    except Exception:
        return []
    try:
        root = ET.fromstring(xml_content)
    except Exception:
        return []
    texts = []
    for elem in root.iter():
        if elem.text and elem.text.strip():
            texts.append(elem.text.strip())
    if texts:
        return [{"title": "Content", "body": "\n".join(texts)}]
    return []


def _parse_zip_xml_table(file_bytes: bytes, xml_path: str = "content.xml") -> list[dict]:
    """Parse table data from a zip-based XML spreadsheet (ODS)."""
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
        if xml_path not in zf.namelist():
            zf.close()
            return []
        xml_content = zf.read(xml_path)
        zf.close()
    except Exception:
        return []
    try:
        root = ET.fromstring(xml_content)
    except Exception:
        return []
    sections = []
    table_ns = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"
    text_ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    for table_elem in root.iter(f"{{{table_ns}}}table"):
        table_name = table_elem.get(f"{{{table_ns}}}name", "Sheet")
        rows = []
        for row_elem in table_elem.iter(f"{{{table_ns}}}table-row"):
            cells = []
            for cell_elem in row_elem.iter(f"{{{table_ns}}}table-cell"):
                cell_texts = []
                for p in cell_elem.iter(f"{{{text_ns}}}p"):
                    if p.text:
                        cell_texts.append(p.text)
                cells.append(" ".join(cell_texts))
            if any(c.strip() for c in cells):
                rows.append(cells)
        if not rows:
            continue
        headers = rows[0]
        data_rows = rows[1:]
        md_rows = [[str(c) if c else "" for c in row] for row in data_rows]
        table_md = _table_to_markdown(headers, md_rows)
        summary = _generate_table_summary(table_name, headers, len(data_rows), table_name)
        sections.append({"title": f"{table_name}", "body": summary, "type": "table_summary"})
        sections.append({"title": f"{table_name} - Data", "body": table_md, "type": "table_data"})
    return sections


# ═══════════════════════════════════════════════════════════════════════════
# Parser Dispatcher
# ═══════════════════════════════════════════════════════════════════════════

PARSER_REGISTRY = {
    # E-book formats
    ".epub": ("EPUB", parse_epub),
    ".mobi": ("MOBI", parse_mobi),
    ".azw": ("MOBI", parse_mobi),
    ".azw3": ("MOBI", parse_mobi),
    ".fb2": ("FB2", parse_fb2),
    ".djvu": ("DjVu", parse_djvu),
    # Office documents
    ".docx": ("Word", parse_docx),
    ".odt": ("ODT", parse_odt),
    ".pptx": ("PowerPoint", parse_pptx),
    ".odp": ("Presentation", parse_odp),
    ".ods": ("Spreadsheet", parse_ods),
    ".xlsx": ("Excel", parse_xlsx),
    ".xls": ("Excel", parse_xlsx),
    ".csv": ("CSV", parse_csv),
    # Markup
    ".md": ("Markdown", parse_markdown),
    ".markdown": ("Markdown", parse_markdown),
    ".html": ("HTML", parse_html),
    ".htm": ("HTML", parse_html),
    ".xml": ("XML", parse_xml),
    ".svg": ("SVG", parse_svg),
    ".yaml": ("YAML", parse_yaml),
    ".yml": ("YAML", parse_yaml),
    ".json": ("JSON", parse_json),
    # Text
    ".txt": ("Text", parse_text),
    ".log": ("Text", parse_text),
    ".ini": ("Text", parse_text),
    ".cfg": ("Text", parse_text),
    ".conf": ("Text", parse_text),
    ".rtf": ("RTF", parse_rtf),
    # LaTeX
    ".tex": ("LaTeX", parse_latex),
    ".sty": ("LaTeX", parse_latex),
    ".cls": ("LaTeX", parse_latex),
    ".bib": ("LaTeX", parse_latex),
}


def parse_file(file_bytes: bytes, filename: str, doc_title: str = "") -> dict:
    """
    Parse a file by its extension.

    Args:
        file_bytes: Raw file content
        filename: Original filename (used to detect type)
        doc_title: Optional document title (used for table summaries)

    Returns:
        dict with either:
          {"sections": [...]} for structured formats
          {"text": "..."} for plain-text formats
    """
    ext = os.path.splitext(filename.lower())[1]
    entry = PARSER_REGISTRY.get(ext)
    if not entry:
        return {"text": "", "sections": []}

    file_type, parser_func = entry

    # Table-aware parsers need doc_title
    if file_type in ("CSV", "Excel"):
        sections = parser_func(file_bytes, doc_title or os.path.splitext(filename)[0])
    elif file_type == "Markdown":
        text = file_bytes.decode("utf-8", errors="ignore")
        sections = parser_func(text)
    else:
        sections = parser_func(file_bytes)

    if sections:
        # Check if any section has type "table_data" or "table_summary"
        has_table = any(s.get("type") in ("table_data", "table_summary") for s in sections)
        if has_table:
            return {"sections": sections, "file_type": file_type}
        return {"sections": sections, "file_type": file_type}
    else:
        # Fallback: return raw text
        try:
            text = file_bytes.decode("utf-8", errors="ignore")
        except Exception:
            text = ""
        return {"text": text, "file_type": file_type}


def extract_preview(file_bytes: bytes, filename: str, max_chars: int = 1000) -> str:
    """
    Extract a text preview from a file for auto-tagging classification.
    Uses the same parsers but only reads enough to classify the content.
    """
    result = parse_file(file_bytes, filename)
    if result.get("sections"):
        for section in result["sections"]:
            body = section.get("body", "") or ""
            if body.strip():
                return body[:max_chars]
    text = result.get("text", "")
    if text.strip():
        return text[:max_chars]
    return f"[{result.get('file_type', 'Unknown')} file: {filename}]"


def detect_file_type(filename: str) -> str:
    """Detect file type from filename extension."""
    ext = os.path.splitext(filename.lower())[1]
    entry = PARSER_REGISTRY.get(ext)
    if entry:
        return entry[0]
    return ""
